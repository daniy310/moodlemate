import os
import time
import requests
import pdfplumber
import datetime
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from pymongo import MongoClient
from telegram import Update
from telegram.ext import CallbackContext
from pptx import Presentation
from bs4 import BeautifulSoup
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options

from telegram_agent.mongo_handler import MongoHandler

# Configure headless browser
options = webdriver.ChromeOptions()
#options.add_argument("--headless=new")
driver = webdriver.Chrome(options=options)

# MongoDB connection
students_collection = MongoHandler("students")
lectures_collection = MongoHandler("lectures")

async def login_to_sols(email: str, password: str, update: Update, context: CallbackContext):
    """Logs into SOLSS and handles OTP verification."""
    try:
        driver.get("https://solss.uow.edu.au/sid/sols_login_ctl.login_app")
        time.sleep(7)

        # Enter email & continue
        driver.find_element(By.NAME, "loginfmt").send_keys("ds374@uowmail.edu.au")
        driver.find_element(By.ID, "idSIButton9").click()
        time.sleep(5)

        # Enter password & continue
        driver.find_element(By.NAME, "passwd").send_keys("fyrT7wvY")
        driver.find_element(By.ID, "idSIButton9").click()
        time.sleep(5)

        # Find OTP code and send to user (Replace with Telegram message logic)
        otp_code = driver.find_element(By.ID, "idRichContext_DisplaySign").text.strip()
        await update.message.reply_text(f"Your OTP number is {otp_code}. Please enter it on your phone.")
        time.sleep(10)

        # Click 'Yes' to stay signed in
        driver.find_element(By.ID, "idSIButton9").click()
        time.sleep(5)

    except Exception as e:
        print(f"❌ Login Error: {e}")
        return None

def get_mobile_timetable():
    timetable = {
        "Monday": [],
        "Tuesday": [],
        "Wednesday": [],
        "Thursday": [],
        "Friday": []
    }

    # Resize the window to simulate a mobile screen (e.g., 375px width)
    driver.set_window_size(375, 800)  # Simulate mobile view

    # Give the page some time to adjust to the mobile layout
    driver.implicitly_wait(3)

    # Find the timetable in mobile view (which is now a list of items)
    try:
        # Find all the <ul> elements (timetable days)
        ul_elements = driver.find_elements(By.XPATH, "//ul[@class='list-group']")
        print(f"✅ Found {len(ul_elements)} timetable sections in mobile view.")
    except Exception as e:
        print(f"❌ Error: Could not find mobile timetable - {e}")
        return timetable

    current_day = None
    for ul in ul_elements:
        li_elements = ul.find_elements(By.TAG_NAME, "li")

        for li in li_elements:
            # Check for day headers safely
            day_headers = li.find_elements(By.XPATH, ".//h4[@class='list-group-item-heading']")
            if day_headers:
                possible_day = day_headers[0].text.strip()
                if possible_day in timetable:  # Only update if it's a valid weekday
                    current_day = possible_day
                    print(f"📅 Found day: {current_day}")

            # Ensure `current_day` is set before processing courses
            if current_day is None:
                continue  # Skip processing if no valid day is found yet

            # Find enrolled courses safely
            course_headers = li.find_elements(By.XPATH, ".//h4[contains(text(), 'Enrolled')]")
            if course_headers:
                course_id = course_headers[0].text.strip().split('-')[-1].strip()

                # Extract course details safely
                course_details_elements = li.find_elements(By.XPATH, ".//p[@class='list-group-item-text']")
                course_details = course_details_elements[0].text.strip() if course_details_elements else "No details available"

                # Split course details into structured information
                course_details_lines = course_details.split("\n")
                if len(course_details_lines) >= 3:
                    course_type = course_details_lines[0].split(':')[0].strip()
                    course_timing = course_details_lines[1].split(',')[-1].strip()
                    course_start_time, course_end_time = course_timing.split('-')
                    course_location = course_details_lines[2].split(':')[-1].strip()

                    print(f"📌 Found course: {course_id} - {course_type} \n Starting at : {course_start_time} - {course_end_time} \n Location: {course_location}")

                    # Add course to timetable
                    timetable[current_day].append({
                        "course_id": course_id,
                        "course_type": course_type,
                        "course_start_time": course_start_time.strip(),
                        "course_end_time": course_end_time.strip(),
                        "course_location": course_location
                    })

    return timetable


def scrape_timetable():
    """Scrapes the student's timetable."""
    try:
        # Go to 'My Timetable' page
        driver.find_element(By.XPATH, "//a[h2[text()='My Timetable']]").click()
        time.sleep(5)

        # Extract timetable content
        timetable = get_mobile_timetable(driver)

        # Resize window back to normal & return to dashboard
        driver.set_window_size(1200, 800)
        driver.back()
        time.sleep(3)

        return timetable

    except Exception as e:
        print(f"❌ Error scraping timetable: {e}")
        return None


def open_moodle():
    """Navigates to Moodle from the SOLSS dashboard."""
    try:
        time.sleep(5)
        moodle_button = driver.find_element(By.XPATH, "//a[h2[text()='Learning Platform (Moodle)']]")
        moodle_button.click()
        time.sleep(10)

        # Switch to the new Moodle tab
        driver.switch_to.window(driver.window_handles[-1])

    except Exception as e:
        print(f"❌ Error opening Moodle: {e}")


# ----- 4️⃣ SCRAPE LECTURE PDFs -----
def scrape_lectures(moodle_url, driver_path):
    """Scrapes Moodle for lecture materials and saves to DB."""
    print(f"📘 Scraping lectures from: {moodle_url}")

    # Set up Selenium driver
    options = Options()
    options.add_argument("--headless")  # Run in headless mode
    options.add_argument("--disable-gpu")
    service = Service(driver_path)
    driver = webdriver.Chrome(service=service, options=options)

    try:
        driver.get(moodle_url)
        soup = BeautifulSoup(driver.page_source, "html.parser")

        # Find all lecture resource links
        lecture_links = soup.find_all("a", class_="rui-course-card")
        print(f"🔍 Found {len(lecture_links)} lecture links.")

        for link in lecture_links:
            lecture_title = link.text.strip()
            lecture_url = link["href"]

            print(f"📥 Downloading: {lecture_title}")
            file_path = download_presentation(lecture_url)

            if file_path:
                lecture_text = process_downloaded_file(file_path)
                if lecture_text:
                    save_lecture_to_db(lecture_title, lecture_url, lecture_text)

    finally:
        print()
        # driver.quit()

def download_presentation(url):
    """Downloads a presentation and saves with the correct extension."""
    try:
        response = requests.get(url, stream=True)

        if response.status_code != 200:
            print(f"❌ Failed to download file: {response.status_code}")
            return None

        # Detect file type
        content_type = response.headers.get("Content-Type", "").lower()
        print(f"📄 Detected Content-Type: {content_type}")

        # Determine file extension
        if "presentationml.presentation" in content_type:
            file_extension = "pptx"
        elif "pdf" in content_type:
            file_extension = "pdf"
        else:
            print("⚠️ Unknown file type. Saving as .bin")
            file_extension = "bin"

        # Save file
        file_path = f"temp_lecture.{file_extension}"
        with open(file_path, "wb") as file:
            for chunk in response.iter_content(chunk_size=1024):
                file.write(chunk)

        print(f"✅ File saved: {file_path}")
        return file_path

    except Exception as e:
        print(f"❌ Error downloading file: {e}")
        return None

def process_downloaded_file(file_path):
    """Extracts text from a downloaded lecture file."""
    if file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    elif file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    else:
        print("⚠️ Unsupported file type")
        return None

def extract_text_from_pptx(pptx_path):
    """Extracts text from a PowerPoint file."""
    try:
        if not os.path.exists(pptx_path):
            print(f"❌ PPTX file does not exist: {pptx_path}")
            return None

        prs = Presentation(pptx_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())

        os.remove(pptx_path)  # Clean up after extraction
        return "\n".join(text)

    except Exception as e:
        print(f"❌ Error extracting text from PPTX: {e}")
        return None

def extract_text_from_pdf(pdf_path):
    """Extracts text from a PDF file."""
    try:
        if not os.path.exists(pdf_path):
            print(f"❌ PDF file does not exist: {pdf_path}")
            return None

        text = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text())

        os.remove(pdf_path)  # Clean up after extraction
        return "\n".join(text)

    except Exception as e:
        print(f"❌ Error extracting text from PDF: {e}")
        return None



def save_lecture_to_db(title, url, text):
    """Saves extracted lecture text to MongoDB."""
    lecture_data = {
        "lecture_title": title,
        "lecture_url": url,
        "lecture_text": text
    }
    lectures_collection.save(lecture_data, unique_key={"lecture_url": url})
    print(f"✅ Lecture saved: {title}")




# ----- 7️⃣ MAIN FUNCTION -----
async def scrape_user_data(user_email, user_password, user_id, update, context):
    """Complete scraping process: Login → Timetable → Moodle → Lectures."""

    try:
        # Login to SOLS
        await login_to_sols(user_email, user_password, update, context)

        # Open Moodle
        open_moodle()

        timetable = scrape_timetable()

        return timetable


        # # Scrape lectures and save to DB
        # lectures = scrape_lectures(user_id)
        #
        # # Save lectures to MongoDB
        # lectures_collection.insert_many(lectures)

        print("✅ Scraping complete!")

    except Exception as e:
        print(f"❌ Error during scraping: {e}")

    finally:
        driver.quit()  # Ensure the driver is closed properly
